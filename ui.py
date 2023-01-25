import sys

from PyQt5.QtCore import QThread, QObject, pyqtSignal, pyqtSlot
from PyQt5.QtWidgets import QApplication, QWidget, QGridLayout, QListWidget, QVBoxLayout, QLineEdit, QTextEdit, QLabel, \
    QPushButton, QScrollArea
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
import sys
from selenium import webdriver
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.common.by import By
from selenium.webdriver.support.select import Select
from selenium.webdriver.chrome.options import Options
from webdriver_manager.chrome import ChromeDriverManager
from ui import *

def trap_exc_during_debug(*args):
    # when app raises uncaught exception, print info
    print(args)


# install exception hook: without this, uncaught exception would cause application to exit
sys.excepthook = trap_exc_during_debug

class App(QWidget):
    sig_start = pyqtSignal()  # needed only due to PyCharm debugger bug (!)
    sig_abort_workers = pyqtSignal()
    def __init__(self):
        super().__init__()
        self.files = []
        self.list = None
        self.layout = None
        self.button = None
        self.scroll = None
        self.initUI()

    def initUI(self):
        self.layout = QVBoxLayout()
        self.setWindowTitle("최새암 바보")
        self.move(300, 300)
        self.resize(1000, 600)

        lay = QVBoxLayout()
        self.button = QPushButton("확인")
        self.button.clicked.connect(self.clicked)
        self.list = QWidget()
        self.list.setLayout(self.layout)

        self.setAcceptDrops(True)

        lay.addWidget(self.list)
        lay.addWidget(self.button)

        self.setLayout(lay)
        self.show()

    def text_on_shape(self, shape, input_text, color, font_size=32, bold=True, clear=True):
        # shape 내 텍스트 프레임 선택하기 & 기존 값 삭제하기
        text_frame = shape.text_frame
        if clear:
            text_frame.clear()

        # 문단 선택하기
        p = text_frame.paragraphs[0]

        # 정렬 설정 : 중간정렬
        p.alighnment = PP_PARAGRAPH_ALIGNMENT.CENTER

        # 텍스트 입력 / 폰트 지정
        run = p.add_run()
        run.text = input_text
        font = run.font
        font.size = Pt(font_size)
        font.bold = bold
        font.color.rgb = color
        font.name = None  # 지정하지 않으면 기본 글씨체로  #  'Arial'

    def make_ppt(self, list):
        bible_abrv = {
            "창": "창세기", "출": "출애굽기", "레": "레위기",
            "민": "민수기", "신": "신명기", "수": "여호수아",
            "삿": "사사기", "룻": "룻기", "삼상": "사무엘상",
            "삼하": "사무엘하", "왕상": "열왕기상", "왕하": "열왕기하",
            "대상": "역대상", "대하": "역대하", "스": "에스라",
            "느": "느헤미야", "에": "에스더", "욥": "욥기",
            "시": "시편", "잠": "잠언", "전": "전도서",
            "아": "아가", "사": "이사야", "렘": "예레미야",
            "애": "예레미야애가", "겔": "에스겔", "단": "다니엘",
            "호": "호세아", "욜": "요엘", "암": "아모스",
            "옵": "오바댜", "욘": "요나", "미": "미가",
            "나": "나훔", "합": "하박국", "습": "스바냐",
            "학": "학개", "슥": "스가랴", "말": "말라기",
            "마": "마태복음", "막": "마가복음", "눅": "누가복음",
            "요": "요한복음", "행": "사도행전", "롬": "로마서",
            "고전": "고린도전서", "고후": "고린도후서", "갈": "갈라디아서",
            "엡": "에베소서", "빌": "빌립보서", "골": "골로새서",
            "살전": "데살로니가전서", "살후": "데살로니가후서", "딤전": "디모데전서",
            "딤후": "디모데후서", "딛": "디도서", "몬": "빌레몬서",
            "히": "히브리서", "약": "야고보서", "벧전": "베드로전서",
            "벧후": "베드로후서", "요일": "요한일서", "요이": "요한이서",
            "요삼": "요한삼서", "유": "유다서", "계": "요한계시록",
        }
        url = "http://www.holybible.or.kr/"
        TABLE_SIZE = 5

        # 크롤링
        options = Options()
        options.add_argument('--blink-settings=imageEnabled=false')
        driver = webdriver.Chrome(options=options, service=Service(ChromeDriverManager().install()))

        for i in range(len(list)):
            # inputs split
            arg = list[i]
            seperated = arg.split('-')
            year = seperated[0][:4].strip()
            month = seperated[0][4:6].strip()
            day = seperated[0][6:].strip()
            contents = seperated[1].strip()
            temp = seperated[2].strip().split(",")  # 요6_61~63, 계2_2~3 fix here

            what = []
            for i in range(len(temp)):
                if "~" in temp[i]:
                    what.append([bible_abrv[temp[i][0]], temp[i][1], temp[i].split("_")[1].split("~")[0],
                                 temp[i].split("_")[1].split("~")[1]])
                else:
                    what.append([bible_abrv[temp[i][0]], temp[i][1], temp[i].split("_")[1]])

            # words = bible_abrv[temp[0].strip()]
            # chapter = temp[1].strip()
            # clause = temp[2].strip()

            pastor = seperated[3].strip().split(".")[0]

            # pptx config
            fpath = "./template.pptx"
            prs = Presentation(fpath)

            # slide 1
            slide = prs.slides[0]
            shape_list = slide.shapes
            shape_index = {}
            for i, shape in enumerate(shape_list):
                shape_index[shape.name] = i

            self.text_on_shape(shape_list[shape_index["title"]], pastor + " 설교", RGBColor(255, 255, 255), font_size=40)
            self.text_on_shape(shape_list[shape_index["date"]], year + "년 " + month + "월 " + day + "일",
                               RGBColor(181, 251, 5))
            self.text_on_shape(shape_list[shape_index["content"]], "'" + contents + "'", RGBColor(255, 255, 255))
            result = ""
            for i in range(len(what)):
                if len(what[i]) == 3:
                    result += (what[i][0] + " " + what[i][1] + "장 " + what[i][2] + "절 ")
                else:
                    result += (what[i][0] + " " + what[i][1] + "장 " + what[i][2] + "~" + what[i][3] + "절 ")

            self.text_on_shape(shape_list[shape_index["words"]], result, RGBColor(181, 251, 5))

            # slide 2
            slide = prs.slides[1]
            shape_list = slide.shapes
            shape_index = {}
            for i, shape in enumerate(shape_list):
                shape_index[shape.name] = i

            self.text_on_shape(shape_list[shape_index["title"]], "'" + contents + "' ", RGBColor(255, 255, 255),
                               font_size=28)
            self.text_on_shape(shape_list[shape_index["title"]], result, RGBColor(181, 251, 5), font_size=28,
                               clear=False)

            gae = []
            saenew = []
            # 개역개정
            for i in range(len(what)):
                driver.get(url)
                select = Select(driver.find_elements(By.NAME, "VR")[0])
                select.select_by_value('GAE')
                search = driver.find_element(By.NAME, "QR")
                search.send_keys(what[i][0] + what[i][1] + "장")
                submit = driver.find_element(By.CSS_SELECTOR,
                                             "body > table:nth-child(3) > tbody > tr > td > table:nth-child(2) > tbody > tr:nth-child(2) > td:nth-child(2) > table > tbody > tr > td > font.tk2 > input[type=submit]:nth-child(2)")
                submit.click()

                if len(what[i]) == 3:
                    clauses = [what[i][2]]
                else:
                    clauses = range(int(what[i][2]), int(what[i][3]) + 1)

                for clause in clauses:
                    index = ((int(clause) - 1) // TABLE_SIZE)
                    list_index = ((int(clause) - 1) % TABLE_SIZE)
                    num = index * TABLE_SIZE + 1
                    clause_id = "b_" + "{0:03d}".format(num)

                    ol = driver.find_element(By.ID, clause_id)
                    gae.append(ol.find_elements(By.TAG_NAME, "li")[list_index].find_element(By.TAG_NAME, "font").text)
                # 새번역
                driver.get(url)
                select = Select(driver.find_elements(By.NAME, "VR")[0])
                select.select_by_value('SAENEW')
                search = driver.find_element(By.NAME, "QR")
                search.send_keys(what[i][0] + what[i][1] + "장")
                submit = driver.find_element(By.CSS_SELECTOR,
                                             "body > table:nth-child(3) > tbody > tr > td > table:nth-child(2) > tbody > tr:nth-child(2) > td:nth-child(2) > table > tbody > tr > td > font.tk2 > input[type=submit]:nth-child(2)")
                submit.click()

                for clause in clauses:
                    index = ((int(clause) - 1) // TABLE_SIZE)
                    list_index = ((int(clause) - 1) % TABLE_SIZE)
                    num = index * TABLE_SIZE + 1
                    clause_id = "b_" + "{0:03d}".format(num)
                    ol = driver.find_element(By.ID, clause_id)
                    saenew.append(
                        ol.find_elements(By.TAG_NAME, "li")[list_index].find_element(By.TAG_NAME, "font").text)

            left = 2.7
            top = 3.41
            width = 30
            height = 2.86
            for k in range(len(saenew)):
                tb1 = slide.shapes.add_textbox(Cm(left), Cm(top + (2 * k * height)), Cm(width), Cm(height))
                tb2 = slide.shapes.add_textbox(Cm(left), Cm(top + (2 * k * height) + height), Cm(width), Cm(height))
                tf1 = tb1.text_frame
                tf2 = tb2.text_frame
                tf1.word_wrap = True
                tf2.word_wrap = True
                self.text_on_shape(tb1, saenew[k], RGBColor(255, 255, 255), font_size=28)
                self.text_on_shape(tb2, gae[k], RGBColor(204, 255, 204), font_size=28)

            # slide 3
            slide = prs.slides[2]
            shape_list = slide.shapes
            shape_index = {}
            for i, shape in enumerate(shape_list):
                shape_index[shape.name] = i

            self.text_on_shape(shape_list[shape_index["title"]], "'" + contents + "' ", RGBColor(255, 255, 255),
                               font_size=28)
            self.text_on_shape(shape_list[shape_index["title"]], result,
                               RGBColor(181, 251, 5), font_size=28, clear=False)

            for k in range(len(saenew)):
                tb1 = slide.shapes.add_textbox(Cm(left), Cm(top + (2 * k * height)), Cm(width), Cm(height))
                tb2 = slide.shapes.add_textbox(Cm(left), Cm(top + (2 * k * height) + height), Cm(width), Cm(height))
                tf1 = tb1.text_frame
                tf2 = tb2.text_frame
                tf1.word_wrap = True
                tf2.word_wrap = True
                self.text_on_shape(tb1, saenew[k], RGBColor(255, 255, 255), font_size=28)
                self.text_on_shape(tb2, gae[k], RGBColor(204, 255, 204), font_size=28)

            prs.save(seperated[0] + seperated[1] + seperated[2] + seperated[3] + ".pptx")

        driver.quit()

    def clicked(self):
        self.make_ppt(self.files)
        self.files.clear()
        for i in range(self.layout.count()):
            self.layout.itemAt(i).widget().deleteLater()

        print("working")

    def dragEnterEvent(self, event):
        if event.mimeData().hasUrls():
            event.accept()
        else:
            event.ignore()

    def dropEvent(self, event):
        for u in event.mimeData().urls():
            self.files.append(u.toLocalFile().split("/")[-1])
            item = QLabel(u.toLocalFile().split("/")[-1], self)
            item.setStyleSheet("border: 1px solid black")
            self.layout.addWidget(item)

