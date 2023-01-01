import collections
import collections.abc
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
import sys
from selenium import webdriver
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.common.by import By
from selenium.webdriver.support.select import Select
from selenium.webdriver.chrome.options import Options
from webdriver_manager.chrome import ChromeDriverManager

TABLE_SIZE = 5

def text_on_shape(shape, input_text, color, font_size=32, bold=True, clear=True):
    # shape 내 텍스트 프레임 선택하기 & 기존 값 삭제하기
    text_frame = shape.text_frame
    if clear:
        text_frame.clear()

    # 문단 선택하기
    p = text_frame.paragraphs[0]

    # 정렬 설정 : 중간정렬
    p.alighnment = PP_PARAGRAPH_ALIGNMENT.CENTER

    # 텍스트 입력 / 폰트 지정
    run = p.add_run()
    run.text = input_text
    font = run.font
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color
    font.name = None  # 지정하지 않으면 기본 글씨체로  #  'Arial'


if __name__ == '__main__':

    bible_abrv = {
        "창": "창세기", "출": "출애굽기", "레": "레위기",
        "민": "민수기", "신": "신명기", "수": "여호수아",
        "삿": "사사기", "룻": "룻기", "삼상": "사무엘상",
        "삼하": "사무엘하", "왕상": "열왕기상", "왕하": "열왕기하",
        "대상": "역대상", "대하": "역대하", "스": "에스라",
        "느": "느헤미야", "에": "에스더", "욥": "욥기",
        "시": "시편", "잠": "잠언", "전": "전도서",
        "아": "아가", "사": "이사야", "렘": "예레미야",
        "애": "예레미야애가", "겔": "에스겔", "단": "다니엘",
        "호": "호세아", "욜": "요엘", "암": "아모스",
        "옵": "오바댜", "욘": "요나", "미": "미가",
        "나": "나훔", "합": "하박국", "습": "스바냐",
        "학": "학개", "슥": "스가랴", "말": "말라기",
        "마": "마태복음", "막": "마가복음", "눅": "누가복음",
        "요": "요한복음", "행": "사도행전", "롬": "로마서",
        "고전": "고린도전서", "고후": "고린도후서", "갈": "갈라디아서",
        "엡": "에베소서", "빌": "빌립보서", "골": "골로새서",
        "살전": "데살로니가전서", "살후": "데살로니가후서", "딤전": "디모데전서",
        "딤후": "디모데후서", "딛": "디도서", "몬": "빌레몬서",
        "히": "히브리서", "약": "야고보서", "벧전": "베드로전서",
        "벧후": "베드로후서", "요일": "요한일서", "요이": "요한이서",
        "요삼": "요한삼서", "유": "유다서", "계": "요한계시록",
    }
    url = "http://www.holybible.or.kr/"

    # 크롤링
    options = Options()
    options.add_argument('--blink-settings=imageEnabled=false')
    driver = webdriver.Chrome(options=options, service=Service(ChromeDriverManager().install()))

    for i in range(len(sys.argv)):
        if i == 0:
            continue
        # inputs split
        arg = sys.argv[i]
        seperated = arg.split('/')
        year = seperated[0][:4].strip()
        month = seperated[0][4:6].strip()
        day = seperated[0][6:].strip()
        contents = seperated[1].strip()
        temp = seperated[2].strip().split(" ")
        words = bible_abrv[temp[0].strip()]
        chapter = temp[1].strip()
        clause = temp[2].strip()
        pastor = seperated[3].strip()


        # pptx config
        fpath = "./template.pptx"
        prs = Presentation(fpath)


        # slide 1
        slide = prs.slides[0]
        shape_list = slide.shapes
        shape_index = {}
        for i, shape in enumerate(shape_list):
            shape_index[shape.name] = i

        text_on_shape(shape_list[shape_index["title"]], pastor + " 설교", RGBColor(255,255,255), font_size=40)
        text_on_shape(shape_list[shape_index["date"]], year+"년 " + month + "월 " + day + "일", RGBColor(181,251,5))
        text_on_shape(shape_list[shape_index["content"]], "'" + contents + "'", RGBColor(255,255,255))
        text_on_shape(shape_list[shape_index["words"]], words + " " + chapter + "장 " + clause + "절",RGBColor(181,251,5))

        #slide 2
        slide = prs.slides[1]
        shape_list = slide.shapes
        shape_index = {}
        for i, shape in enumerate(shape_list):
            shape_index[shape.name] = i

        text_on_shape(shape_list[shape_index["title"]], "'" + contents + "' ", RGBColor(255,255,255), font_size=28)
        text_on_shape(shape_list[shape_index["title"]], words + " " + chapter + "장 " + clause + "절", RGBColor(181,251,5), font_size=28, clear=False)


        #개역개정
        driver.get(url)
        select = Select(driver.find_elements(By.NAME, "VR")[0])
        select.select_by_value('GAE')
        search = driver.find_element(By.NAME, "QR")
        search.send_keys(words+chapter+"장")
        submit = driver.find_element(By.CSS_SELECTOR, "body > table:nth-child(3) > tbody > tr > td > table:nth-child(2) > tbody > tr:nth-child(2) > td:nth-child(2) > table > tbody > tr > td > font.tk2 > input[type=submit]:nth-child(2)")
        submit.click()

        index = ((int(clause) - 1) // TABLE_SIZE)
        list_index = ((int(clause) - 1) % TABLE_SIZE)
        num = index * TABLE_SIZE + 1
        clause_id = "b_" + "{0:03d}".format(num)
        ol = driver.find_element(By.ID, clause_id)

        gae = ol.find_elements(By.TAG_NAME, "li")[list_index].find_element(By.TAG_NAME, "font").text

        #새번역
        driver.get(url)
        select = Select(driver.find_elements(By.NAME, "VR")[0])
        select.select_by_value('SAENEW')
        search = driver.find_element(By.NAME, "QR")
        search.send_keys(words + chapter + "장")
        submit = driver.find_element(By.CSS_SELECTOR,
                                     "body > table:nth-child(3) > tbody > tr > td > table:nth-child(2) > tbody > tr:nth-child(2) > td:nth-child(2) > table > tbody > tr > td > font.tk2 > input[type=submit]:nth-child(2)")
        submit.click()

        index = ((int(clause) - 1) // TABLE_SIZE)
        list_index = ((int(clause) - 1) % TABLE_SIZE)
        num = index * TABLE_SIZE + 1
        clause_id = "b_" + "{0:03d}".format(num)
        ol = driver.find_element(By.ID, clause_id)

        saenew = ol.find_elements(By.TAG_NAME, "li")[list_index].find_element(By.TAG_NAME, "font").text

        text_on_shape(shape_list[shape_index["new_translation"]], saenew, RGBColor(255, 255, 255))
        text_on_shape(shape_list[shape_index["revision"]], gae, RGBColor(204, 255, 204))


        #slide 3
        slide = prs.slides[2]
        shape_list = slide.shapes
        shape_index = {}
        for i, shape in enumerate(shape_list):
            shape_index[shape.name] = i

        text_on_shape(shape_list[shape_index["title"]], "'" + contents + "' ", RGBColor(255, 255, 255), font_size=28)
        text_on_shape(shape_list[shape_index["title"]], words + " " + chapter + "장 " + clause + "절",
                      RGBColor(181, 251, 5), font_size=28, clear=False)
        text_on_shape(shape_list[shape_index["new_translation"]], saenew, RGBColor(255, 255, 255))
        text_on_shape(shape_list[shape_index["revision"]], gae, RGBColor(204, 255, 204))

        prs.save(seperated[0] + " " + seperated[1] + " " + seperated[2] + " " + seperated[3] + ".pptx")

    driver.quit()
